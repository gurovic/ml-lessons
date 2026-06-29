"""
Сборщик презентации из JSON-слайдов.

Использование: python agents/pptx_builder.py <lesson_dir>
"""

import json
import sys
from datetime import date
from pathlib import Path


def read_file(path: Path) -> str:
    with open(path, "r", encoding="utf-8") as f:
        return f.read()


def load_config(root_dir: Path) -> dict:
    """Загружает project_config.json из корня проекта, если есть."""
    cfg_path = root_dir / "project_config.json"
    if cfg_path.exists():
        return json.loads(read_file(cfg_path))
    return {}


def load_slides(slides_dir: Path) -> list[dict]:
    from slide_utils import list_slide_files

    return [json.loads(read_file(f)) for f in list_slide_files(slides_dir)]


def _image_aspect(img_path: Path) -> float:
    from PIL import Image

    with Image.open(img_path) as im:
        w, h = im.size
    return w / h if h else 1.0


def _visual_size_hint(visual_obj: dict) -> str:
    """Размер из JSON: large | default | compact."""
    size = visual_obj.get("size", "default")
    if size not in ("large", "default", "compact"):
        return "default"
    return size


BODY_LEFT = 0.7
BODY_WIDTH_FULL = 12.3
BULLET_TOP = 1.5
BULLET_COL_GAP = 0.2
# Две колонки буллетов на всю ширину слайда (только без картинки)
FULL_COL_W = (BODY_WIDTH_FULL - BULLET_COL_GAP) / 2
# Левая колонка текста при картинке справа
BODY_WIDTH_VIS = 6.85
BULLET_COL_W = 3.4
IMG_COL_LEFT = 7.95
IMG_COL_W = 5.5
VISUAL_TOP = BULLET_TOP  # правая колонка: верх картинки = верх буллетов
TEXT_IMG_GAP = 0.18
BULLET_CODE_GAP = 0.38
BULLET_CODE_GAP_COL = 0.55  # слайд: текст слева + картинка справа + код внизу слева
VISUAL_STACK_GAP = 0.16
IMG_VERTICAL_GAP = 0.34  # ≈ 1 строка буллета между картинками в колонке
STACKED_IMG_GAP = 0.22
WIDE_BAND_LEFT = 0.7
WIDE_BAND_W = 7.9
MIN_WIDE_BAND_H = 0.85
MAX_WIDE_BAND_H = 2.45
TWO_COL_MIN_BULLETS = 3
BULLET_CHAR = "\u2022"  # • — единый маркер для всех слайдов
SLIDE_WIDTH = 13.333
SLIDE_HEIGHT = 7.5
SLIDE_RIGHT_MARGIN = 0.5
CODE_BOTTOM_MARGIN = 0.12  # блок кода прижат к низу слайда


def _max_image_right(slide_width: float = SLIDE_WIDTH) -> float:
    return slide_width - SLIDE_RIGHT_MARGIN


def _clamp_image_width(
    left: float, width: float, *, max_right: float | None = None
) -> float:
    limit = max_right if max_right is not None else _max_image_right()
    return min(width, max(limit - left, 0.0))


def _image_column_dims(*, max_right: float | None = None) -> tuple[float, float]:
    col_w = _clamp_image_width(IMG_COL_LEFT, IMG_COL_W, max_right=max_right)
    return IMG_COL_LEFT, col_w


def _wide_band_dims(*, max_right: float | None = None) -> tuple[float, float]:
    wide_w = _clamp_image_width(WIDE_BAND_LEFT, WIDE_BAND_W, max_right=max_right)
    return WIDE_BAND_LEFT, wide_w


def _format_bullet_text(text: str) -> str:
    """Префикс маркера для pptx; в JSON bullets[] — без «•»."""
    stripped = text.lstrip()
    for prefix in (f"{BULLET_CHAR} ", f"{BULLET_CHAR}\t", "- ", "* "):
        if stripped.startswith(prefix):
            stripped = stripped[len(prefix) :].lstrip()
            return f"{BULLET_CHAR} {stripped}"
    if stripped.startswith(BULLET_CHAR):
        stripped = stripped[len(BULLET_CHAR) :].lstrip()
    return f"{BULLET_CHAR} {stripped}"


def _split_bullets(bullets: list[str]) -> tuple[list[str], list[str]]:
    mid = (len(bullets) + 1) // 2
    return bullets[:mid], bullets[mid:]


def _use_two_column_bullets(bullets: list[str], has_visuals: bool) -> bool:
    """С картинкой: одна колонка слева. Без картинки: две колонки при ≥5 буллетов."""
    if has_visuals:
        return False
    if len(bullets) < TWO_COL_MIN_BULLETS:
        return False
    return len(bullets) >= 5


def _would_use_wide_band(aspects: list[float], size_hints: list[str]) -> bool:
    if len(aspects) != 1:
        return False
    hint = size_hints[0] if size_hints else "default"
    return hint == "large"


def _estimate_bullet_height(
    bullets: list[str],
    *,
    two_col: bool,
    col_w: float,
) -> float:
    """Оценка высоты блока буллетов с учётом переноса строк."""
    if not bullets:
        return 0.0
    chars_per_line = max(18, int(col_w * 7.2))
    if two_col:
        left, right = _split_bullets(bullets)
        cols = [left, right]
    else:
        cols = [bullets]
    max_lines = 0
    for col in cols:
        lines = sum(max(1, (len(b) + chars_per_line - 1) // chars_per_line) for b in col)
        max_lines = max(max_lines, lines)
    return 0.30 * max_lines + 0.35


def _compute_slide_layout(
    bullets: list[str],
    code_examples: list[dict],
    visuals: list[dict],
    assets_dir: Path,
    *,
    has_visuals: bool,
    has_link: bool,
    estimate_code_height_inches,
) -> dict:
    content_bottom = 6.0 if has_link else 6.35
    bullet_top = BULLET_TOP
    stacked = False
    two_col = _use_two_column_bullets(bullets, has_visuals)
    has_code = bool(code_examples)
    body_width = BODY_WIDTH_VIS if has_visuals else BODY_WIDTH_FULL
    code_w = body_width - 0.1
    code_h = (
        estimate_code_height_inches(code_examples, code_width=code_w)
        if has_code
        else 0.0
    )
    min_bullet_h = 0.9 if bullets else 0.0

    bullet_col_w = FULL_COL_W if two_col else body_width

    aspects: list[float] = []
    size_hints: list[str] = []
    if has_visuals:
        for visual_obj in visuals:
            output = visual_obj.get("output", "")
            if not output:
                continue
            img_path = assets_dir / output
            if img_path.exists():
                aspects.append(_image_aspect(img_path))
                size_hints.append(_visual_size_hint(visual_obj))

    if two_col and bullets:
        left_b, right_b = _split_bullets(bullets)
        n_rows = max(len(left_b), len(right_b))
    else:
        n_rows = len(bullets)

    est_bullet_h = _estimate_bullet_height(bullets, two_col=two_col, col_w=bullet_col_w)

    bullet_h = 0.0
    bullet_bottom = bullet_top

    column_with_code = has_code and has_visuals and not stacked
    bullet_clipped = False

    if has_code:
        code_top = SLIDE_HEIGHT - CODE_BOTTOM_MARGIN - code_h
    else:
        code_top = None

    if bullets:
        if stacked and code_top is not None:
            max_h = code_top - bullet_top - STACKED_IMG_GAP - MIN_WIDE_BAND_H
            bullet_h = max(min_bullet_h, min(max_h, est_bullet_h))
        elif stacked:
            max_h = content_bottom - bullet_top - STACKED_IMG_GAP - MIN_WIDE_BAND_H
            bullet_h = max(min_bullet_h, min(max_h, est_bullet_h))
        elif column_with_code and code_top is not None:
            gap = BULLET_CODE_GAP_COL
            max_h = code_top - bullet_top - gap
            bullet_h = max(min_bullet_h, min(max_h, est_bullet_h))
            bullet_clipped = est_bullet_h > max_h + 0.05
        elif code_top is not None:
            max_h = code_top - bullet_top - BULLET_CODE_GAP
            bullet_h = max(min_bullet_h, min(max_h, est_bullet_h))
            bullet_clipped = est_bullet_h > max_h + 0.05
        else:
            bullet_h = min(5.2, max(2.0, est_bullet_h))
        bullet_bottom = bullet_top + bullet_h

    visual_top = (bullet_bottom + STACKED_IMG_GAP) if stacked else BULLET_TOP
    visual_bottom = content_bottom - 0.08

    if has_code and code_top is not None and column_with_code:
        visual_bottom = min(visual_bottom, code_top - TEXT_IMG_GAP)

    wide_band = (not stacked) and _would_use_wide_band(aspects, size_hints)
    band_top = band_h = None

    if wide_band:
        band_top = (bullet_bottom + TEXT_IMG_GAP) if bullets else BULLET_TOP
        band_h = min(MAX_WIDE_BAND_H, visual_bottom - band_top)
        if band_h < MIN_WIDE_BAND_H:
            wide_band = False
            band_top = band_h = None

    return {
        "body_left": BODY_LEFT,
        "body_width": body_width,
        "bullet_col_w": bullet_col_w,
        "bullet_top": bullet_top,
        "bullet_h": bullet_h,
        "bullet_bottom": bullet_bottom,
        "two_col": two_col,
        "stacked": stacked,
        "wide_band": wide_band,
        "band_top": band_top,
        "band_h": band_h,
        "visual_top": visual_top,
        "visual_bottom": visual_bottom,
        "code_top": code_top,
        "content_bottom": content_bottom,
        "bullet_clipped": bullet_clipped,
    }


def _picture_fit_rect(left, top, box_w, box_h, aspect: float, *, top_align: bool = True):
    box_aspect = box_w / box_h if box_h else aspect
    if aspect >= box_aspect:
        width = box_w
        height = width / aspect
    else:
        height = box_h
        width = height * aspect
    x = left + (box_w - width) / 2
    y = top if top_align else top + (box_h - height) / 2
    return (x, y, width, height)


def _rects_intersect(a, b) -> bool:
    ax, ay, aw, ah = a
    bx, by, bw, bh = b
    return ax < bx + bw and ax + aw > bx and ay < by + bh and ay + ah > by


def _bullet_text_rects(layout: dict) -> list[tuple[float, float, float, float]]:
    if layout["bullet_h"] <= 0:
        return []
    left = layout["body_left"]
    top = layout["bullet_top"]
    h = layout["bullet_h"]
    col_w = layout.get("bullet_col_w", BULLET_COL_W)
    if layout["two_col"]:
        return [
            (left, top, col_w, h),
            (left + col_w + BULLET_COL_GAP, top, col_w, h),
        ]
    return [(left, top, layout["body_width"], h)]


def _visual_layout_slots(
    n: int,
    col_left: float,
    col_w: float,
    top: float,
    bottom: float,
    *,
    aspects: list[float] | None = None,
    size_hints: list[str] | None = None,
    layout: dict | None = None,
):
    """Координаты ячеек (left, top, width, height) в дюймах для 1–4 картинок."""
    gap = VISUAL_STACK_GAP if layout.get("stacked") else IMG_VERTICAL_GAP
    layout = layout or {}

    if layout.get("stacked"):
        band_left = BODY_LEFT
        band_w = BODY_WIDTH_FULL
        vtop = layout.get("visual_top", top)
        vbottom = layout.get("visual_bottom", bottom)
        h_avail = max(0.5, vbottom - vtop)
        if n <= 0:
            return []
        if n == 1:
            return [(band_left, vtop, band_w, h_avail)]
        if n == 2:
            slot_h = (h_avail - gap) / 2
            return [
                (band_left, vtop, band_w, slot_h),
                (band_left, vtop + slot_h + gap, band_w, slot_h),
            ]
        slot_w = (band_w - gap) / 2
        slot_h = (h_avail - gap) / 2
        if n == 3:
            return [
                (band_left, vtop, slot_w, slot_h),
                (band_left + slot_w + gap, vtop, slot_w, slot_h),
                (band_left, vtop + slot_h + gap, band_w, slot_h),
            ]
        return [
            (band_left, vtop, slot_w, slot_h),
            (band_left + slot_w + gap, vtop, slot_w, slot_h),
            (band_left, vtop + slot_h + gap, slot_w, slot_h),
            (band_left + slot_w + gap, vtop + slot_h + gap, slot_w, slot_h),
        ]

    h_avail = bottom - top
    aspects = aspects or []
    size_hints = size_hints or []

    if n <= 0:
        return []
    if n == 1:
        use_wide = layout.get("wide_band") and layout.get("band_top") is not None
        if use_wide and layout.get("band_h"):
            wide_left, wide_w = _wide_band_dims()
            return [(wide_left, layout["band_top"], wide_w, layout["band_h"])]
        slot_bottom = layout.get("visual_bottom", bottom)
        slot_h = max(0.5, slot_bottom - top)
        return [(col_left, top, col_w, slot_h)]
    if n == 2:
        slot_h = (h_avail - gap) / 2
        return [
            (col_left, top, col_w, slot_h),
            (col_left, top + slot_h + gap, col_w, slot_h),
        ]

    slot_w = (col_w - gap) / 2
    slot_h = (h_avail - gap) / 2
    if n == 3:
        return [
            (col_left, top, slot_w, slot_h),
            (col_left + slot_w + gap, top, slot_w, slot_h),
            (col_left, top + slot_h + gap, col_w, slot_h),
        ]
    return [
        (col_left, top, slot_w, slot_h),
        (col_left + slot_w + gap, top, slot_w, slot_h),
        (col_left, top + slot_h + gap, slot_w, slot_h),
        (col_left + slot_w + gap, top + slot_h + gap, slot_w, slot_h),
    ]


def _add_picture_fit(slide, img_path: Path, left, top, box_w, box_h):
    from pptx.util import Inches

    x, y, width, height = _picture_fit_rect(left, top, box_w, box_h, _image_aspect(img_path), top_align=True)
    slide.shapes.add_picture(str(img_path), Inches(x), Inches(y), width=Inches(width))


def _add_hyperlink_run(paragraph, text: str, url: str, font_size, link_rgb):
    from pptx.dml.color import RGBColor

    run = paragraph.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = RGBColor(*link_rgb)
    run.hyperlink.address = url


def _render_references_slide(
    slide,
    slide_data: dict,
    assets_dir: Path,
    slide_num: int,
    *,
    set_paragraph_content,
    Pt,
    Inches,
    body_rgb,
    link_rgb,
):
    """Слайд type=references: литература слева, QR для Colab справа."""
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    from references_utils import format_paper_bullet, generate_qr_png

    refs = slide_data.get("references", [])
    papers = [r for r in refs if r.get("kind") == "paper"]
    colabs = [r for r in refs if r.get("kind") == "colab"]

    body_left, body_top = 0.7, 1.45
    body_width = 7.2
    y_cursor = body_top

    def add_section_heading(text: str):
        nonlocal y_cursor
        box = slide.shapes.add_textbox(
            Inches(body_left), Inches(y_cursor), Inches(body_width), Inches(0.35)
        )
        p = box.text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*body_rgb)
        y_cursor += 0.38

    def add_text_line(height: float = 0.55):
        nonlocal y_cursor
        box = slide.shapes.add_textbox(
            Inches(body_left), Inches(y_cursor), Inches(body_width), Inches(height)
        )
        tf = box.text_frame
        tf.word_wrap = True
        y_cursor += height
        return tf.paragraphs[0]

    if papers:
        add_section_heading("Литература")
        for ref in papers:
            p = add_text_line(0.62)
            p.space_after = Pt(4)
            bullet_text = format_paper_bullet(ref)
            url = ref.get("url")
            if url:
                for w in set_paragraph_content(
                    p, f"{BULLET_CHAR} ", font_size=Pt(16), color=body_rgb
                ):
                    print(f"  [warn] Слайд {slide_num}, литература: {w}")
                _add_hyperlink_run(p, bullet_text, url, Pt(16), link_rgb)
            else:
                for w in set_paragraph_content(
                    p, _format_bullet_text(bullet_text), font_size=Pt(16), color=body_rgb
                ):
                    print(f"  [warn] Слайд {slide_num}, литература: {w}")

    if colabs:
        add_section_heading("Практика в Colab")
        for entry in colabs:
            url = entry.get("url", "")
            label = entry.get("label") or entry.get("title", "Colab")
            p = add_text_line(0.5)
            for w in set_paragraph_content(
                p, _format_bullet_text(f"{label}: "), font_size=Pt(15), color=body_rgb
            ):
                print(f"  [warn] Слайд {slide_num}, Colab: {w}")
            if url:
                _add_hyperlink_run(p, url, url, Pt(14), link_rgb)

    bullets = slide_data.get("bullets", [])
    if bullets:
        y_cursor += 0.1
        for bullet in bullets:
            p = add_text_line(0.45)
            for w in set_paragraph_content(
                p, _format_bullet_text(bullet), font_size=Pt(15), color=body_rgb
            ):
                print(f"  [warn] Слайд {slide_num}, буллет: {w}")

    qr_size = 1.15
    qr_gap = 0.25
    qr_left = 9.0
    qr_top_start = 1.6
    for i, entry in enumerate(colabs):
        url = entry.get("url", "")
        if not url:
            continue
        label = entry.get("label") or entry.get("title", "Colab")
        top = qr_top_start + i * (qr_size + qr_gap + 0.35)
        try:
            qr_path = assets_dir / f"_qr_ref_{slide_num}_{i}.png"
            generate_qr_png(url, qr_path)
            slide.shapes.add_picture(
                str(qr_path), Inches(qr_left), Inches(top), width=Inches(qr_size)
            )
            qr_path.unlink(missing_ok=True)
        except ImportError:
            print("  [warn] Установи qrcode: pip install qrcode[pil]")
            break

        cap = slide.shapes.add_textbox(
            Inches(qr_left - 0.1),
            Inches(top + qr_size + 0.02),
            Inches(qr_size + 0.4),
            Inches(0.3),
        )
        cp = cap.text_frame.paragraphs[0]
        cp.text = label
        cp.font.size = Pt(10)
        cp.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
        cp.alignment = PP_ALIGN.CENTER


def _get_python_highlighter():
    """Pygments lexer + style для подсветки Python; None при отсутствии pygments."""
    try:
        from pygments.lexers import PythonLexer
        from pygments.styles import get_style_by_name

        return PythonLexer(stripnl=False), get_style_by_name("friendly")
    except ImportError:
        return None, None


def _rgb_from_hex(hex_color: str):
    from pptx.dml.color import RGBColor

    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _append_highlighted_code_line(
    paragraph,
    line: str,
    *,
    lexer,
    style,
    code_font: str,
    code_size,
    default_rgb,
):
    from pptx.dml.color import RGBColor

    if lexer is None or style is None:
        run = paragraph.add_run()
        run.text = line if line else " "
        run.font.name = code_font
        run.font.size = code_size
        run.font.color.rgb = RGBColor(*default_rgb)
        return

    from pygments import lex

    for token_type, text in lex(line, lexer):
        if not text:
            continue
        run = paragraph.add_run()
        run.text = text
        run.font.name = code_font
        run.font.size = code_size
        token_style = style.style_for_token(token_type)
        color = token_style.get("color")
        if color:
            run.font.color.rgb = _rgb_from_hex(color)
        else:
            run.font.color.rgb = RGBColor(*default_rgb)
        if token_style.get("bold"):
            run.font.bold = True
        if token_style.get("italic"):
            run.font.italic = True


def _render_code_examples(
    slide,
    code_examples: list[dict],
    *,
    left: float,
    top: float,
    width: float,
    Inches,
    body_rgb,
):
    """Моноширинный блок(и) кода: серая рамка + по одному textbox на строку."""
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
    from pptx.util import Pt

    from slide_code_utils import (
        CODE_BLOCK_GAP,
        CODE_BLOCK_TAIL,
        CODE_CAPTION_H,
        CODE_FONT_PT,
        CODE_LINE_H,
        CODE_PAD_BOTTOM,
        CODE_PAD_TOP,
        CODE_PAD_X,
        apply_tight_code_paragraph,
        max_code_line_length_for_width,
        normalize_code_lines,
    )

    if not code_examples:
        return top

    lexer, style = _get_python_highlighter()
    if lexer is None:
        print("  [warn] Pygments не установлен — код без подсветки (pip install pygments)")
    code_font = "Consolas"
    code_size = Pt(CODE_FONT_PT)
    caption_size = Pt(10)
    y = top
    inner_w = width - 2 * CODE_PAD_X
    max_len = max_code_line_length_for_width(width)

    for block in code_examples:
        source = block.get("source", "")
        if not source:
            continue
        lines = normalize_code_lines(source, max_len=max_len)
        n_lines = len(lines)
        block_h = CODE_PAD_TOP + CODE_LINE_H * n_lines + CODE_PAD_BOTTOM

        bg = slide.shapes.add_shape(
            1,
            Inches(left),
            Inches(y),
            Inches(width),
            Inches(block_h),
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
        bg.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)

        for j, line in enumerate(lines):
            line_box = slide.shapes.add_textbox(
                Inches(left + CODE_PAD_X),
                Inches(y + CODE_PAD_TOP + j * CODE_LINE_H),
                Inches(inner_w),
                Inches(CODE_LINE_H),
            )
            line_box.fill.background()
            line_box.line.fill.background()
            tf = line_box.text_frame
            tf.word_wrap = False
            tf.auto_size = MSO_AUTO_SIZE.NONE
            tf.vertical_anchor = MSO_ANCHOR.TOP
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
            p = tf.paragraphs[0]
            apply_tight_code_paragraph(p)
            _append_highlighted_code_line(
                p,
                line,
                lexer=lexer,
                style=style,
                code_font=code_font,
                code_size=code_size,
                default_rgb=body_rgb,
            )

        y += block_h + CODE_BLOCK_GAP

        caption = block.get("caption")
        if caption:
            cap_box = slide.shapes.add_textbox(
                Inches(left),
                Inches(y),
                Inches(width),
                Inches(0.22),
            )
            cp = cap_box.text_frame.paragraphs[0]
            cp.text = caption
            cp.font.size = caption_size
            cp.font.italic = True
            cp.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
            cp.alignment = PP_ALIGN.LEFT
            y += CODE_CAPTION_H

        y += CODE_BLOCK_TAIL

    return y


def _render_bullets(
    slide,
    bullets: list[str],
    layout: dict,
    *,
    set_paragraph_content,
    Pt,
    Inches,
    body_rgb,
    slide_num: int,
    auto_fit: bool = False,
):
    from pptx.enum.text import MSO_AUTO_SIZE

    if not bullets:
        return

    left = layout["body_left"]
    top = layout["bullet_top"]
    height = layout["bullet_h"]

    def fill_column(col_left: float, col_width: float, col_bullets: list[str], label: str):
        box = slide.shapes.add_textbox(
            Inches(col_left), Inches(top), Inches(col_width), Inches(height)
        )
        tf = box.text_frame
        tf.word_wrap = True
        if auto_fit:
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        for j, bullet in enumerate(col_bullets):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.space_after = Pt(8)
            for w in set_paragraph_content(
                p, _format_bullet_text(bullet), font_size=Pt(20), color=body_rgb
            ):
                print(f"  [warn] Слайд {slide_num}, {label} {j + 1}: {w}")

    if layout["two_col"]:
        left_bullets, right_bullets = _split_bullets(bullets)
        col_w = layout.get("bullet_col_w", BULLET_COL_W)
        fill_column(left, col_w, left_bullets, "буллет")
        fill_column(left + col_w + BULLET_COL_GAP, col_w, right_bullets, "буллет (кол.2)")
    else:
        fill_column(left, layout["body_width"], bullets, "буллет")


def _collect_visual_rects(
    visuals: list[dict],
    assets_dir: Path,
    layout: dict,
) -> list[tuple[float, float, float, float]]:
    items: list[tuple[Path, dict]] = []
    for visual_obj in visuals:
        output = visual_obj.get("output", "")
        if not output:
            continue
        img_path = assets_dir / output
        if img_path.exists():
            items.append((img_path, visual_obj))

    if not items:
        return []

    paths = [p for p, _ in items]
    aspects = [_image_aspect(p) for p in paths]
    size_hints = [_visual_size_hint(v) for _, v in items]
    col_left, col_w = _image_column_dims()
    vtop = layout.get("visual_top", BULLET_TOP)
    vbottom = layout.get("visual_bottom", layout["content_bottom"])
    slots = _visual_layout_slots(
        len(paths),
        col_left,
        col_w,
        vtop,
        vbottom,
        aspects=aspects,
        size_hints=size_hints,
        layout=layout,
    )
    rects = []
    for img_path, (left, slot_top, w, h) in zip(paths, slots):
        rects.append(
            _picture_fit_rect(left, slot_top, w, h, _image_aspect(img_path), top_align=True)
        )
    return rects


def _verify_slide_layouts(
    slides: list[dict],
    assets_dir: Path,
    *,
    estimate_code_height_inches,
) -> list[str]:
    errors: list[str] = []
    for i, slide_data in enumerate(slides, 1):
        if slide_data.get("type") == "references":
            continue
        visuals = slide_data.get("visuals", [])
        has_visuals = any(v.get("output") for v in visuals)
        if not has_visuals:
            continue

        layout = _compute_slide_layout(
            slide_data.get("bullets", []),
            slide_data.get("code_examples", []),
            visuals,
            assets_dir,
            has_visuals=True,
            has_link=bool(slide_data.get("link", {}).get("url")),
            estimate_code_height_inches=estimate_code_height_inches,
        )
        text_rects = _bullet_text_rects(layout)
        image_rects = _collect_visual_rects(visuals, assets_dir, layout)
        max_right = _max_image_right()
        for t_rect in text_rects:
            for img_rect in image_rects:
                if _rects_intersect(t_rect, img_rect):
                    errors.append(
                        f"Слайд {i}: пересечение текста {t_rect} и картинки {img_rect}"
                    )
        for img_rect in image_rects:
            x, _y, w, _h = img_rect
            right = x + w
            if right > max_right + 1e-3:
                errors.append(
                    f"Слайд {i}: картинка выходит за правый край "
                    f"(right={right:.3f}\" > {max_right:.3f}\") rect={img_rect}"
                )
    return errors


def _audit_image_overflow_slides(
    slides: list[dict],
    assets_dir: Path,
    *,
    estimate_code_height_inches,
) -> list[dict]:
    """Слайды, у которых старые размеры колонки выходили за правый край."""
    max_right = _max_image_right()
    old_col_right = IMG_COL_LEFT + IMG_COL_W
    col_left, col_w = _image_column_dims()
    wide_left, wide_w = _wide_band_dims()
    fixed: list[dict] = []

    for i, slide_data in enumerate(slides, 1):
        if slide_data.get("type") == "references":
            continue
        visuals = slide_data.get("visuals", [])
        if not any(v.get("output") for v in visuals):
            continue

        layout = _compute_slide_layout(
            slide_data.get("bullets", []),
            slide_data.get("code_examples", []),
            visuals,
            assets_dir,
            has_visuals=True,
            has_link=bool(slide_data.get("link", {}).get("url")),
            estimate_code_height_inches=estimate_code_height_inches,
        )
        if layout.get("wide_band"):
            old_right = WIDE_BAND_LEFT + 12.0
            new_dims = f"wide left={wide_left}\", w={wide_w:.3f}\""
            overflow = old_right > max_right + 1e-3
        else:
            old_right = old_col_right
            new_dims = f"col left={col_left}\", w={col_w:.3f}\""
            overflow = old_col_right > max_right + 1e-3

        if overflow:
            fixed.append(
                {
                    "slide": i,
                    "layout": "wide_band" if layout.get("wide_band") else "column",
                    "old_right": old_right,
                    "max_right": max_right,
                    "new_dims": new_dims,
                }
            )
    return fixed


def _place_visuals(
    slide,
    visuals: list[dict],
    assets_dir: Path,
    slide_num: int,
    *,
    layout: dict,
):
    items: list[tuple[Path, dict]] = []
    for visual_obj in visuals:
        output = visual_obj.get("output", "")
        if not output:
            continue
        img_path = assets_dir / output
        if img_path.exists():
            items.append((img_path, visual_obj))
        else:
            print(f"  [warn] Слайд {slide_num}: нет файла {img_path.name}")

    if not items:
        return

    paths = [p for p, _ in items]
    size_hints = [_visual_size_hint(v) for _, v in items]
    aspects = [_image_aspect(p) for p in paths]

    col_left, col_w = _image_column_dims()
    vtop = layout.get("visual_top", BULLET_TOP)
    vbottom = layout.get("visual_bottom", layout["content_bottom"])
    slots = _visual_layout_slots(
        len(paths),
        col_left,
        col_w,
        vtop,
        vbottom,
        aspects=aspects,
        size_hints=size_hints,
        layout=layout,
    )
    for (img_path, _), (left, slot_top, w, h) in zip(items, slots):
        try:
            _add_picture_fit(slide, img_path, left, slot_top, w, h)
        except Exception as e:
            print(f"Не удалось вставить {img_path}: {e}")


def build_presentation(slides: list[dict], output_path: Path, assets_dir: Path, lesson_dir: Path):
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
        from rich_text import set_paragraph_content
    except ImportError as e:
        print(f"Установи зависимости: pip install -r requirements.txt ({e})")
        sys.exit(1)

    title_rgb = (0x1A, 0x1A, 0x2E)
    body_rgb = (0x00, 0x00, 0x00)
    link_rgb = (0x33, 0x66, 0xCC)

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH)
    prs.slide_height = Inches(SLIDE_HEIGHT)

    # Титульный слайд
    project_config = load_config(Path(__file__).parent.parent)

    info = {}
    info_path = lesson_dir / "info.json"
    if info_path.exists():
        with open(info_path, "r", encoding="utf-8") as f:
            info = json.load(f)

    lesson_title = info.get("topic", lesson_dir.name)
    author = info.get("author", "") or project_config.get("author", "")
    email = info.get("email", "") or project_config.get("email", "")
    telegram = info.get("telegram", "") or info.get("tg", "") or project_config.get("telegram", "")
    today = date.today().strftime("%d.%m.%Y")

    title_slide = prs.slides.add_slide(prs.slide_layouts[6])

    tb = title_slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = lesson_title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
    p.alignment = PP_ALIGN.CENTER

    if author:
        tb2 = title_slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11), Inches(0.6))
        p2 = tb2.text_frame.paragraphs[0]
        p2.text = f"Автор: {author}"
        p2.font.size = Pt(20)
        p2.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
        p2.alignment = PP_ALIGN.CENTER

    contacts = " | ".join(filter(None, [f"Email: {email}" if email else "", f"Telegram: {telegram}" if telegram else ""]))
    if contacts:
        tb3 = title_slide.shapes.add_textbox(Inches(1), Inches(4.7), Inches(11), Inches(0.6))
        p3 = tb3.text_frame.paragraphs[0]
        p3.text = contacts
        p3.font.size = Pt(18)
        p3.font.color.rgb = RGBColor(0x77, 0x77, 0x77)
        p3.alignment = PP_ALIGN.CENTER

    tb4 = title_slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11), Inches(0.5))
    p4 = tb4.text_frame.paragraphs[0]
    p4.text = today
    p4.font.size = Pt(16)
    p4.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    p4.alignment = PP_ALIGN.CENTER

    # Слайды
    for i, slide_data in enumerate(slides, 1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        if slide_data.get("formula"):
            print(f"  [warn] Слайд {i}: поле formula устарело — используйте $...$ в тексте (docs/formulas.md)")

        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(1))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        for w in set_paragraph_content(
            p,
            slide_data.get("title", f"Слайд {i}"),
            font_size=Pt(32),
            bold=True,
            color=title_rgb,
        ):
            print(f"  [warn] Слайд {i}, заголовок: {w}")

        visuals = slide_data.get("visuals", [])
        has_visuals = any(v.get("output") for v in visuals)
        is_references = slide_data.get("type") == "references"

        if is_references:
            _render_references_slide(
                slide,
                slide_data,
                assets_dir,
                i,
                set_paragraph_content=set_paragraph_content,
                Pt=Pt,
                Inches=Inches,
                body_rgb=body_rgb,
                link_rgb=link_rgb,
            )
        else:
            from slide_code_utils import estimate_code_height_inches

            code_examples = slide_data.get("code_examples", [])
            bullets = slide_data.get("bullets", [])
            has_link = bool(slide_data.get("link", {}).get("url"))

            layout = _compute_slide_layout(
                bullets,
                code_examples,
                visuals,
                assets_dir,
                has_visuals=has_visuals,
                has_link=has_link,
                estimate_code_height_inches=estimate_code_height_inches,
            )

            if bullets:
                _render_bullets(
                    slide,
                    bullets,
                    layout,
                    set_paragraph_content=set_paragraph_content,
                    Pt=Pt,
                    Inches=Inches,
                    body_rgb=body_rgb,
                    slide_num=i,
                    auto_fit=layout.get("bullet_clipped", False),
                )

            if code_examples and layout["code_top"] is not None:
                _render_code_examples(
                    slide,
                    code_examples,
                    left=layout["body_left"],
                    top=layout["code_top"],
                    width=layout["body_width"] - 0.1,
                    Inches=Inches,
                    body_rgb=body_rgb,
                )

            _place_visuals(slide, visuals, assets_dir, i, layout=layout)

            # Ссылка и QR-код
            link = slide_data.get("link")
            if link and link.get("url"):
                url = link["url"]
                label = link.get("label", url)

                link_box = slide.shapes.add_textbox(Inches(0.7), Inches(6.5), Inches(6), Inches(0.5))
                tf_link = link_box.text_frame
                tf_link.word_wrap = True
                p_link = tf_link.paragraphs[0]
                for w in set_paragraph_content(
                    p_link, f"{label}: ", font_size=Pt(14), color=link_rgb
                ):
                    print(f"  [warn] Слайд {i}, ссылка: {w}")
                _add_hyperlink_run(p_link, url, url, Pt(14), link_rgb)

                try:
                    from references_utils import generate_qr_png

                    qr_path = assets_dir / "_qr_temp.png"
                    generate_qr_png(url, qr_path, box_size=10)
                    slide.shapes.add_picture(str(qr_path), Inches(11.5), Inches(4.5), width=Inches(1.5))
                    qr_path.unlink(missing_ok=True)
                except ImportError:
                    print("  [warn] Установи qrcode: pip install qrcode[pil]")

        notes = slide_data.get("notes")
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    prs.save(str(output_path))
    print(f"Сохранено: {output_path}")


def main():
    if len(sys.argv) < 2:
        print("Использование: python agents/pptx_builder.py <lesson_dir>")
        sys.exit(1)

    lesson_dir = Path(sys.argv[1])
    if not lesson_dir.exists():
        print(f"Папка не найдена: {lesson_dir}")
        sys.exit(1)

    slides_dir = lesson_dir / "slides_json"
    if not slides_dir.exists():
        print(f"Папка со слайдами не найдена: {slides_dir}")
        sys.exit(1)

    slides = load_slides(slides_dir)
    if not slides:
        print(f"Нет JSON-файлов в {slides_dir}")
        sys.exit(1)

    print(f"Загружено слайдов: {len(slides)}")
    from slide_code_utils import estimate_code_height_inches

    assets_dir = lesson_dir / "assets"
    col_left, col_w = _image_column_dims()
    wide_left, wide_w = _wide_band_dims()
    print("  [layout] column: text left (1 col with visuals), image right")

    overflow_slides = _audit_image_overflow_slides(
        slides, assets_dir, estimate_code_height_inches=estimate_code_height_inches
    )
    if overflow_slides:
        print("  [layout] Исправлен выход за правый край на слайдах:")
        for entry in overflow_slides:
            print(
                f"    слайд {entry['slide']} ({entry['layout']}): "
                f"было right={entry['old_right']:.3f}\", "
                f"теперь {entry['new_dims']}"
            )

    layout_errors = _verify_slide_layouts(
        slides, assets_dir, estimate_code_height_inches=estimate_code_height_inches
    )
    if layout_errors:
        for err in layout_errors:
            print(f"  [layout] {err}")
        print("  [layout] Обнаружены проблемы вёрстки — проверьте pptx_builder.py")
    else:
        visual_slides = sum(
            1
            for s in slides
            if s.get("type") != "references"
            and any(v.get("output") for v in s.get("visuals", []))
        )
        two_col_slides = [
            i + 1
            for i, s in enumerate(slides)
            if s.get("type") != "references"
            and any(v.get("output") for v in s.get("visuals", []))
            and _use_two_column_bullets(s.get("bullets", []), True)
        ]
        print(f"  [layout] OK: {visual_slides} слайдов с иллюстрациями, без пересечений")
        if two_col_slides:
            print(f"  [layout] Двухколоночные буллеты: слайды {two_col_slides}")

    build_presentation(slides, lesson_dir / "presentation.pptx", assets_dir, lesson_dir)
    print("Откройте presentation.pptx для проверки (JSON править не обязательно).")


if __name__ == "__main__":
    main()
